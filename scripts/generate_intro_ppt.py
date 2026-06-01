"""QuickPaste Manager 소개 PPT 생성 (10장, 미니멀)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "QuickPaste_Manager_소개.pptx"

ACCENT = RGBColor(0x1A, 0x5F, 0xB4)
TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _blank_layout(prs: Presentation):
    return prs.slide_layouts[6]


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_para(tf, text: str, *, size: int, bold: bool = False, color=TEXT, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(8)


def _accent_bar(slide, prs: Presentation) -> None:
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def _slide_title(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _accent_bar(slide, prs)
    box = _textbox(slide, Inches(0.9), Inches(2.2), Inches(8.2), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    _set_para(tf, title, size=36, bold=True)
    if subtitle:
        _set_para(tf, subtitle, size=18, color=MUTED)


def _slide_body(prs: Presentation, title: str, bullets: list[str], footer: str = "") -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _accent_bar(slide, prs)
    head = _textbox(slide, Inches(0.9), Inches(0.55), Inches(8.2), Inches(0.8))
    _set_para(head.text_frame, title, size=28, bold=True)

    body = _textbox(slide, Inches(0.9), Inches(1.45), Inches(8.2), Inches(4.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(4)

    if footer:
        foot = _textbox(slide, Inches(0.9), Inches(6.3), Inches(8.2), Inches(0.5))
        _set_para(foot.text_frame, footer, size=14, color=MUTED)


def _slide_flow(prs: Presentation) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    _accent_bar(slide, prs)
    head = _textbox(slide, Inches(0.9), Inches(0.55), Inches(8.2), Inches(0.8))
    _set_para(head.text_frame, "동작 흐름", size=28, bold=True)

    steps = [
        "① 단축키 또는 마우스 가운데 버튼",
        "② 커서 옆 팝업",
        "③ 상용구 선택",
        "④ 클립보드 + 자동 붙여넣기",
    ]
    y = 1.6
    for step in steps:
        box = _textbox(slide, Inches(1.2), Inches(y), Inches(7.5), Inches(0.55))
        tf = box.text_frame
        _set_para(tf, step, size=22)
        y += 0.75
        if step != steps[-1]:
            arrow = _textbox(slide, Inches(1.2), Inches(y - 0.35), Inches(1), Inches(0.3))
            _set_para(arrow.text_frame, "↓", size=18, color=MUTED)
            y += 0.15


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_title(
        prs,
        "QuickPaste Manager",
        "반복 입력을 한 번의 클릭으로  ·  Windows 상용구 유틸리티",
    )

    _slide_body(
        prs,
        "왜 필요한가",
        [
            "같은 문장·이미지를 매일 여러 번 복사·붙여넣기",
            "메모장·채팅·메일을 오가며 찾는 시간 낭비",
            "클립보드 히스토리만으로는 분류·고정이 어렵다",
        ],
    )

    _slide_body(
        prs,
        "무엇을 해주는가",
        [
            "자주 쓰는 텍스트·이미지를 카테고리별로 저장",
            "어떤 프로그램에서든 전역 단축키로 팝업 호출",
            "선택 한 번으로 현재 입력 위치에 붙여넣기",
        ],
    )

    _slide_flow(prs)

    _slide_body(
        prs,
        "활용 ① — 반복 업무 문구",
        [
            "사내 공지·보고 서식, 회의 안내 문구",
            "「고정」으로 Top 5에 상시 노출",
            "태그로 「보고」「공지」 검색",
        ],
        footer="예: 주간 보고 인사말, 퇴근 안내",
    )

    _slide_body(
        prs,
        "활용 ② — 고객 응대·이메일",
        [
            "카테고리: 고객응대 / 이메일 / 일반",
            "FAQ 답변·사과·안내 템플릿을 제목으로 구분",
            "채팅·메일·CRM 입력창 어디서든 동일 단축키",
        ],
    )

    _slide_body(
        prs,
        "활용 ③ — 이미지·연락처",
        [
            "로고·안내 이미지·계좌 캡처를 상용구로 보관",
            "미리보기 확인 후 붙여넣기",
            "ZIP 보내기로 PC 간 데이터 이전",
        ],
    )

    _slide_body(
        prs,
        "핵심만",
        [
            "로컬 저장 — 데이터는 내 PC (%APPDATA%)",
            "트레이 상주 · 단축키 · Top 5 · 검색",
            "설치 파일 1개로 배포 (Setup.exe)",
        ],
    )

    _slide_body(
        prs,
        "시작하기",
        [
            "설치: QuickPasteManager-Setup-0.2.0.exe 실행",
            "트레이 아이콘 → 상용구 관리에서 등록",
            "기본 단축키 Ctrl+Shift+V 로 팝업",
            "도움말·환경설정에서 단축키·백업 조정",
        ],
    )

    slide = prs.slides.add_slide(_blank_layout(prs))
    _accent_bar(slide, prs)
    box = _textbox(slide, Inches(0.9), Inches(2.4), Inches(8.2), Inches(2))
    tf = box.text_frame
    _set_para(tf, "감사합니다", size=32, bold=True, align=PP_ALIGN.CENTER)
    _set_para(
        tf,
        "QuickPaste Manager  ·  반복을 줄이고 본업에 집중하세요",
        size=18,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"생성 완료: {path}")
